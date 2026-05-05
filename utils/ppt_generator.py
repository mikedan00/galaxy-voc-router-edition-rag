"""
utils/ppt_generator.py

VOC 분석 결과를 Microsoft PowerPoint에서 바로 열 수 있는 .pptx로 저장한다.
이 버전은 수작업 XML 패키징 대신 python-pptx를 사용한다.
이전 pure-XML 방식은 LibreOffice에서는 열리지만 Microsoft PowerPoint에서
'읽을 수 없습니다' 오류가 날 수 있어 제거했다.
"""

from __future__ import annotations

import json
import re
import time
from pathlib import Path
from typing import Any, Iterable

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover
    Presentation = None  # type: ignore
    _IMPORT_ERROR = exc
else:
    _IMPORT_ERROR = None


BLUE = RGBColor(20, 40, 160)
SKY = RGBColor(0, 174, 239)
NAVY = RGBColor(17, 24, 39)
GRAY = RGBColor(75, 85, 99)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(229, 231, 235)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(16, 185, 129)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)


def _get(obj: Any, key: str, default: Any = "") -> Any:
    if hasattr(obj, "to_dict"):
        try:
            return obj.to_dict().get(key, default)
        except Exception:
            return default
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)


def _clean(text: Any, max_len: int = 900) -> str:
    text = str(text or "")
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) > max_len:
        return text[: max_len - 3].rstrip() + "..."
    return text


def _list_text(values: Any, limit: int = 6, max_each: int = 120) -> list[str]:
    if not values:
        return []
    if isinstance(values, dict):
        values = [values]
    if not isinstance(values, list):
        values = [values]

    out: list[str] = []
    for item in values[:limit]:
        if isinstance(item, dict):
            title = item.get("title") or item.get("name") or item.get("requirement") or item.get("issue") or item.get("kpi") or ""
            desc = item.get("description") or item.get("summary") or item.get("metric") or item.get("priority") or item.get("acceptance_criteria") or ""
            if title and desc:
                line = f"{title} - {desc}"
            elif title:
                line = title
            else:
                line = json.dumps(item, ensure_ascii=False)
        else:
            line = str(item)
        out.append(_clean(line, max_each))
    return out


def _require_pptx() -> None:
    if Presentation is None:
        raise RuntimeError(
            "PPTX 생성을 위해 python-pptx가 필요합니다. "
            "requirements.txt에 python-pptx를 추가하고 다시 설치/배포하세요. "
            f"원래 import 오류: {_IMPORT_ERROR}"
        )


def _set_font(run, size: int = 18, bold: bool = False, color: RGBColor = NAVY) -> None:
    run.font.name = "맑은 고딕"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(slide, x, y, w, h, text: str = "", font_size: int = 18, bold: bool = False, color: RGBColor = NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, bold, color)
    return box


def _add_card(slide, x, y, w, h, title: str, items: list[str], title_color: RGBColor = NAVY, fill: RGBColor = LIGHT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)

    _add_textbox(slide, x + Inches(0.18), y + Inches(0.13), w - Inches(0.36), Inches(0.38), title, 16, True, title_color)

    body = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.62), w - Inches(0.44), h - Inches(0.78))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    if not items:
        items = ["데이터 없음"]

    for idx, item in enumerate(items[:8]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {_clean(item, 150)}"
        _set_font(run, 12 if len(items) > 5 else 13, False, GRAY)
    return card


def _add_header(slide, title: str, subtitle: str | None = None):
    _add_textbox(slide, Inches(0.45), Inches(0.25), Inches(12.2), Inches(0.45), title, 24, True, BLUE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.48), Inches(0.82), Inches(12.1), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SKY
    line.line.fill.background()
    if subtitle:
        _add_textbox(slide, Inches(0.48), Inches(0.88), Inches(12.0), Inches(0.32), subtitle, 10, False, GRAY)


def _add_footer(slide, text: str = "Galaxy VOC Collector · AI/RAG 기반 자동 생성 보고서"):
    _add_textbox(slide, Inches(0.55), Inches(7.05), Inches(12.0), Inches(0.25), text, 9, False, GRAY, PP_ALIGN.RIGHT)


def _blank_slide(prs):
    # layout 6 is blank in the default PowerPoint-compatible template bundled with python-pptx
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_slide(prs, title: str, subtitle: str, meta: str):
    slide = _blank_slide(prs)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.8), Inches(11.85), Inches(1.5))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.color.rgb = BLUE

    _add_textbox(slide, Inches(1.05), Inches(1.08), Inches(11.1), Inches(0.7), title, 30, True, WHITE, PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.2), Inches(2.75), Inches(10.9), Inches(0.8), subtitle, 20, False, NAVY, PP_ALIGN.CENTER)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(4.35), Inches(5.0), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SKY
    accent.line.fill.background()

    _add_textbox(slide, Inches(1.2), Inches(5.85), Inches(10.9), Inches(0.45), meta, 12, False, GRAY, PP_ALIGN.CENTER)
    return slide


def _metrics_slide(prs, stats: dict, rag_data: dict | None = None):
    slide = _blank_slide(prs)
    _add_header(slide, "VOC 분석 대시보드", "수집량, 감성, 카테고리, 채널 통계를 요약합니다.")

    total = int(stats.get("total", 0) or 0) if isinstance(stats, dict) else 0
    by_sentiment = stats.get("by_sentiment", {}) if isinstance(stats, dict) else {}
    by_category = stats.get("by_category", {}) if isinstance(stats, dict) else {}
    by_source = stats.get("by_source", {}) if isinstance(stats, dict) else {}
    top_cats = sorted(by_category.items(), key=lambda x: x[1], reverse=True)[:5]
    top_sources = sorted(by_source.items(), key=lambda x: x[1], reverse=True)[:5]

    metric_values = [
        ("총 VOC", f"{total:,}", BLUE),
        ("부정 VOC", f"{by_sentiment.get('negative', 0):,}", RED),
        ("긍정 VOC", f"{by_sentiment.get('positive', 0):,}", GREEN),
        ("RAG 청크", f"{(rag_data or {}).get('uploaded_chunks_count', 0):,}", ORANGE),
    ]
    for i, (label, value, color) in enumerate(metric_values):
        x = Inches(0.65 + i * 3.15)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.45), Inches(2.75), Inches(1.15))
        shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = BORDER
        _add_textbox(slide, x + Inches(0.08), Inches(1.62), Inches(2.55), Inches(0.3), label, 12, True, GRAY, PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.08), Inches(1.95), Inches(2.55), Inches(0.5), value, 24, True, color, PP_ALIGN.CENTER)

    _add_card(slide, Inches(0.65), Inches(3.0), Inches(6.0), Inches(3.25), "상위 카테고리", [f"{k}: {v:,}건" for k, v in top_cats] or ["카테고리 통계 없음"])
    _add_card(slide, Inches(6.95), Inches(3.0), Inches(5.75), Inches(3.25), "상위 채널", [f"{k}: {v:,}건" for k, v in top_sources] or ["채널 통계 없음"])
    _add_footer(slide)
    return slide


def _two_column_slide(prs, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str], subtitle: str | None = None):
    slide = _blank_slide(prs)
    _add_header(slide, title, subtitle)
    _add_card(slide, Inches(0.65), Inches(1.35), Inches(6.0), Inches(5.55), left_title, left_items, BLUE, LIGHT)
    _add_card(slide, Inches(6.95), Inches(1.35), Inches(5.75), Inches(5.55), right_title, right_items, NAVY, WHITE)
    _add_footer(slide)
    return slide


def _sample_slide(prs, voc_items: Iterable):
    samples = []
    for i, item in enumerate(list(voc_items)[:8], start=1):
        title = _clean(_get(item, "title", ""), 95)
        source = _get(item, "source", "")
        cat = _get(item, "category", "")
        samples.append(f"{i}. {title} ({source}/{cat})")
    return _two_column_slide(
        prs,
        "대표 VOC 샘플",
        "VOC 예시",
        samples[:4] or ["VOC 샘플 없음"],
        "추가 샘플",
        samples[4:8] or ["추가 샘플 없음"],
        "원본 VOC와 전체 근거는 analysis.json에서 확인할 수 있습니다.",
    )


def _rag_slide(prs, rag_data: dict | None):
    rag_data = rag_data or {}
    hits = rag_data.get("hits") or []
    answer = _clean(rag_data.get("answer", ""), 500)
    hit_lines: list[str] = []
    for h in hits[:6]:
        if isinstance(h, dict):
            hit_lines.append(_clean(h.get("text") or h.get("content") or h.get("title") or json.dumps(h, ensure_ascii=False), 130))
        else:
            hit_lines.append(_clean(h, 130))
    return _two_column_slide(
        prs,
        "RAG 근거 및 답변",
        "검색된 근거",
        hit_lines or ["검색된 RAG 근거 없음"],
        "RAG 답변 요약",
        [answer or "RAG 답변 없음"],
        "업로드 파일과 수집 VOC에서 검색된 근거를 기반으로 생성됩니다.",
    )


def generate_pptx(
    voc_items: Iterable,
    analysis: dict | None,
    srs_text: str = "",
    stats: dict | None = None,
    rag_data: dict | None = None,
    product_name: str = "삼성 갤럭시",
    version: str = "1.0",
    author: str = "제품기획팀",
    output_dir: str = "output",
    model_label: str = "HF Router",
) -> str:
    """PowerPoint에서 직접 열리는 .pptx 보고서를 생성한다."""
    _require_pptx()

    analysis = analysis or {}
    stats = stats or {}

    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    safe_product = re.sub(r"[^0-9A-Za-z가-힣_-]+", "_", product_name).strip("_") or "Galaxy"
    path = out / f"{safe_product}_VOC_Report_v{version}_{time.strftime('%Y%m%d_%H%M%S')}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    meta = f"v{version} · {author} · {time.strftime('%Y-%m-%d')} · {model_label}"
    _title_slide(prs, f"{product_name} VOC 분석 보고서", "VOC 수집, RAG 근거, AI 분석 기반 제품 요구사항 요약", meta)
    _metrics_slide(prs, stats, rag_data)

    _two_column_slide(
        prs,
        "Executive Summary",
        "핵심 요약",
        [_clean(analysis.get("executive_summary", "분석 요약 없음"), 500)],
        "주요 인사이트",
        _list_text(analysis.get("key_insights"), 6, 130),
        "VOC 전체 흐름과 우선 개선 방향을 요약합니다.",
    )
    _two_column_slide(
        prs,
        "핵심 이슈와 문제 진술",
        "문제 진술",
        _list_text(analysis.get("problem_statements"), 6, 130),
        "Critical Issues",
        _list_text(analysis.get("critical_issues"), 6, 130),
    )
    _two_column_slide(
        prs,
        "기능 요구사항",
        "Functional Requirements",
        _list_text(analysis.get("requirements"), 7, 140),
        "수용 기준/메트릭 관점",
        [
            "VOC 근거와 연결된 요구사항 우선순위화",
            "요구사항별 KPI와 Acceptance Criteria 정의",
            "Quick Win과 장기 개선 로드맵 분리",
            "릴리즈 전/후 VOC 변화 추적",
        ],
    )
    _two_column_slide(
        prs,
        "비기능 요구사항과 KPI",
        "Non-functional Requirements",
        _list_text(analysis.get("non_functional_requirements"), 6, 130),
        "KPI",
        _list_text(analysis.get("kpis"), 6, 130),
    )
    _two_column_slide(
        prs,
        "개선 로드맵",
        "Roadmap",
        _list_text(analysis.get("roadmap"), 7, 145),
        "운영 제안",
        [
            "상위 VOC 테마 월간 추적",
            "불만 급증 키워드 알림",
            "릴리즈 전/후 VOC 비교",
            "RAG 근거 기반 요구사항 리뷰",
        ],
    )
    _rag_slide(prs, rag_data)
    _sample_slide(prs, voc_items)

    if srs_text:
        srs_lines = [x.strip("- #") for x in srs_text.splitlines() if x.strip()][:10]
        _two_column_slide(
            prs,
            "SRS 본문 요약",
            "SRS 주요 항목",
            [_clean(x, 130) for x in srs_lines[:5]] or ["SRS 요약 없음"],
            "추가 항목",
            [_clean(x, 130) for x in srs_lines[5:10]] or ["추가 항목 없음"],
        )

    prs.save(path)
    return str(path)
